import zipfile
import os
import re
import shutil
import subprocess
import json

import pdfplumber
import pandas as pd
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfMerger


# =========================
# CONFIGURAÇÕES
# =========================

ZIP_PATH = "input.zip"

BASE_DIR = "workspace"
EXTRACT_DIR = os.path.join(BASE_DIR, "extracted")
RSMA_DIR = os.path.join(BASE_DIR, "filtered_rsma")
CONVERTED_PDF_DIR = os.path.join(BASE_DIR, "converted_pdf")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")

FINAL_PDF = os.path.join(OUTPUT_DIR, "RSMA_final.pdf")
FINAL_ZIP = os.path.join(OUTPUT_DIR, "RSMA_pdfs.zip")

LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"

SUPPORTED_EXTENSIONS = [".pdf", ".xlsx", ".pptx", ".docx", ".doc"]


# =========================
# CONTROLE DE PROGRESSO
# =========================

PROGRESS = {
    "total_files": 0,
    "analyzed_files": 0,
    "rsma_files": 0,
    "converted_files": 0
}


# =========================
# PADRÕES RSMA
# =========================

RSMA_PATTERNS = [
    r"\brsma\b",
    r"rate[\s\-]?splitting",
    r"rate[\s\-]?splitting[\s\-]?multiple[\s\-]?access",
    r"rate[\s\-]?splitting[\s\-]?ma"
]


def is_rsma_text(text: str) -> bool:
    text = text.lower()
    return any(re.search(pattern, text) for pattern in RSMA_PATTERNS)


# =========================
# EXTRAÇÃO ZIP
# =========================

def extract_zip(zip_path, extract_dir):
    if os.path.exists(extract_dir):
        shutil.rmtree(extract_dir)
    os.makedirs(extract_dir, exist_ok=True)

    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)


def extract_nested_zips(base_dir):
    extracted = True
    while extracted:
        extracted = False
        for root, _, files in os.walk(base_dir):
            for file in files:
                if file.lower().endswith(".zip"):
                    zip_path = os.path.join(root, file)
                    extract_path = os.path.splitext(zip_path)[0]
                    try:
                        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                            zip_ref.extractall(extract_path)
                        os.remove(zip_path)
                        extracted = True
                    except Exception:
                        pass


# =========================
# CONTAGEM DE ARQUIVOS
# =========================

def count_supported_files(base_dir):
    count = 0
    for root, _, files in os.walk(base_dir):
        for file in files:
            if os.path.splitext(file)[1].lower() in SUPPORTED_EXTENSIONS:
                count += 1
    return count


# =========================
# EXTRAÇÃO DE TEXTO
# =========================

def extract_text_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                if page.extract_text():
                    text += page.extract_text() + "\n"
    except Exception:
        pass
    return text


def extract_text_xlsx(file_path):
    text = ""
    try:
        xls = pd.ExcelFile(file_path)
        for sheet in xls.sheet_names:
            df = xls.parse(sheet)
            text += df.astype(str).to_string()
    except Exception:
        pass
    return text


def extract_text_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception:
        pass
    return text


def extract_text_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception:
        pass
    return text


def convert_to_pdf_windows(input_file, output_dir):
    os.makedirs(output_dir, exist_ok=True)

    command = [
        LIBREOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_file
    ]

    subprocess.run(command, check=True)

def extract_text_doc_via_pdf(file_path, temp_dir):
    """
    Converte .doc para PDF usando LibreOffice e extrai o texto do PDF.
    """
    os.makedirs(temp_dir, exist_ok=True)

    convert_to_pdf_windows(file_path, temp_dir)

    base_name = os.path.splitext(os.path.basename(file_path))[0]
    pdf_path = os.path.join(temp_dir, base_name + ".pdf")

    return extract_text_pdf(pdf_path)



# =========================
# FILTRAGEM RSMA
# =========================

def filter_rsma_files():
    if os.path.exists(RSMA_DIR):
        shutil.rmtree(RSMA_DIR)
    os.makedirs(RSMA_DIR, exist_ok=True)

    for root, _, files in os.walk(EXTRACT_DIR):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                continue

            PROGRESS["analyzed_files"] += 1
            remaining = PROGRESS["total_files"] - PROGRESS["analyzed_files"]

            print(
                f"🔍 Analisando {PROGRESS['analyzed_files']}/"
                f"{PROGRESS['total_files']} | Restantes: {remaining}"
            )

            file_path = os.path.join(root, file)

            if is_rsma_text(file):
                shutil.copy(file_path, RSMA_DIR)
                PROGRESS["rsma_files"] += 1
                print(f"✅ RSMA identificado (nome do arquivo)")
                continue

            if ext == ".pdf":
                text = extract_text_pdf(file_path)
            elif ext == ".xlsx":
                text = extract_text_xlsx(file_path)
            elif ext == ".pptx":
                text = extract_text_pptx(file_path)
            elif ext == ".docx":
                text = extract_text_docx(file_path)
            elif ext == ".doc":
                text = extract_text_doc_via_pdf(
                    file_path,
                    temp_dir=os.path.join(BASE_DIR, "temp_doc_pdf")
                )
            else:
                continue

            if is_rsma_text(text):
                shutil.copy(file_path, RSMA_DIR)
                PROGRESS["rsma_files"] += 1
                print(f"✅ RSMA identificado (conteúdo)")


# =========================
# CONVERSÃO PARA PDF
# =========================

def convert_to_pdf_windows(input_file, output_dir):
    os.makedirs(output_dir, exist_ok=True)

    command = [
        LIBREOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_file
    ]

    subprocess.run(command, check=True)


def convert_all_to_pdf():
    for file in os.listdir(RSMA_DIR):
        file_path = os.path.join(RSMA_DIR, file)
        ext = os.path.splitext(file)[1].lower()

        if ext == ".pdf":
            shutil.copy(file_path, CONVERTED_PDF_DIR)
            PROGRESS["converted_files"] += 1

        elif ext in [".xlsx", ".pptx", ".docx", ".doc"]:
            convert_to_pdf_windows(file_path, CONVERTED_PDF_DIR)
            PROGRESS["converted_files"] += 1


        print(f"📄 Convertidos: {PROGRESS['converted_files']}")


# =========================
# MERGE E ZIP FINAL
# =========================

def merge_pdfs(output_pdf):
    os.makedirs(os.path.dirname(output_pdf), exist_ok=True)

    merger = PdfMerger()
    for file in sorted(os.listdir(CONVERTED_PDF_DIR)):
        if file.lower().endswith(".pdf"):
            merger.append(os.path.join(CONVERTED_PDF_DIR, file))

    merger.write(output_pdf)
    merger.close()


def zip_final_pdfs(zip_path):
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for file in os.listdir(CONVERTED_PDF_DIR):
            if file.lower().endswith(".pdf"):
                zipf.write(
                    os.path.join(CONVERTED_PDF_DIR, file),
                    arcname=file
                )


# =========================
# MAIN
# =========================

def main():
    print("📦 Extraindo ZIP principal...")
    extract_zip(ZIP_PATH, EXTRACT_DIR)

    print("📦 Extraindo ZIPs internos...")
    extract_nested_zips(EXTRACT_DIR)

    PROGRESS["total_files"] = count_supported_files(EXTRACT_DIR)
    print(f"📊 Total de arquivos a analisar: {PROGRESS['total_files']}")

    print("🔎 Iniciando filtragem RSMA...")
    filter_rsma_files()

    print("📄 Convertendo para PDF...")
    os.makedirs(CONVERTED_PDF_DIR, exist_ok=True)
    convert_all_to_pdf()

    print("🧩 Unificando PDFs...")
    merge_pdfs(FINAL_PDF)

    print("🗜️ Criando ZIP final...")
    zip_final_pdfs(FINAL_ZIP)

    with open(os.path.join(OUTPUT_DIR, "progress.json"), "w") as f:
        json.dump(PROGRESS, f, indent=2)

    print("\n✅ Processo concluído com sucesso!")
    print(f"📄 PDF final: {FINAL_PDF}")
    print(f"🗜️ ZIP final: {FINAL_ZIP}")


if __name__ == "__main__":
    main()
